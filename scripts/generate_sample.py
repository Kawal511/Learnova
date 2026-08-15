"""
Generate the sample PPTX fixtures used by the test suite and verify scripts.

Run:  python scripts/generate_sample.py

Writes two files into tests/fixtures/:
  * sample_test_presentation.pptx — text only
  * sample_with_images.pptx       — pictures on specific slides, used to test
                                    that images anchor to the right section
"""

import io
import pathlib

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

FIXTURES = pathlib.Path(__file__).resolve().parent.parent / "tests" / "fixtures"
DEFAULT_OUTPUT = FIXTURES / "sample_test_presentation.pptx"
IMAGE_OUTPUT = FIXTURES / "sample_with_images.pptx"


def create_sample_presentation(output_path=None) -> str:
    """Build the sample deck and save it. Returns the path written."""
    target = pathlib.Path(output_path) if output_path else DEFAULT_OUTPUT
    target.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = "Artificial Intelligence in Education"
    subtitle1.text = "Transforming Outdated Courseware into Engaging Learning Modules"

    # Slide 2: Minimal text list
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "Key Challenges in Higher Education"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Traditional presentations suffer from heavy text density and passive student engagement."
    p1 = tf2.add_paragraph()
    p1.text = "Lack of interactive assessment checkpoints during lectures reduces retention rates."
    p2 = tf2.add_paragraph()
    p2.text = "Manual redesign of slides requires hours of design expertise per lecture."

    # Slide 3: Process workflow (flowchart candidate)
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "Learnova Automated Processing Pipeline Workflow"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Step 1: Upload raw PPTX or PDF textbook."
    p31 = tf3.add_paragraph()
    p31.text = "Step 2: Parse document layout, tables, SmartArt, and embedded diagram images."
    p32 = tf3.add_paragraph()
    p32.text = "Step 3: Route content through AI Layout Router to classify dynamic visual structures."
    p33 = tf3.add_paragraph()
    p33.text = "Step 4: Generate interactive quizzes and export animated PPTX and HTML web decks."

    # Slide 4: Metric callout candidate
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "Student Engagement Benchmark Metrics"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Studies show 84% improvement in concept retention when interactive checkpoints and dynamic visual layouts are integrated into lecture slides."

    # Slide 5: Card Grid candidate
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "Four Pillars of Learnova Architecture"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "1. AI Visual Engine: Converts plain text to flowcharts and structured cards."
    p51 = tf5.add_paragraph()
    p51.text = "2. Vision OCR Integration: Reads embedded diagrams and infographic text."
    p52 = tf5.add_paragraph()
    p52.text = "3. Interleaved Quizzes: Promotes active recall with checkpoint questions."
    p53 = tf5.add_paragraph()
    p53.text = "4. Multi-Format Export: Native PPTX presentation + 60fps web deck."

    prs.save(str(target))
    return str(target)


def _swatch(color: tuple, label: str) -> io.BytesIO:
    """A synthetic diagram image, so the fixture needs no binary assets in git."""
    img = Image.new("RGB", (480, 320), color)
    draw = ImageDraw.Draw(img)
    draw.rectangle([10, 10, 470, 310], outline=(0, 0, 0), width=6)
    draw.text((30, 150), label, fill=(0, 0, 0))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def create_image_presentation(output_path=None) -> str:
    """
    Build a deck where only *some* slides carry a picture.

    The gaps are the point: they prove images anchor to the section that
    discusses them rather than to a positional index.
    """
    target = pathlib.Path(output_path) if output_path else IMAGE_OUTPUT
    target.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)

    specs = [
        ("Photosynthesis Overview",
         "Plants convert light into chemical energy.", None),
        ("Chloroplast Structure Diagram",
         "The thylakoid membrane stacks form grana.", ((120, 200, 120), "CHLOROPLAST")),
        ("Water Cycle Stages",
         "Evaporation, condensation, precipitation, collection.", None),
        ("Carbon Cycle Diagram",
         "Carbon moves between atmosphere, biosphere and oceans.",
         ((150, 150, 220), "CARBON CYCLE")),
    ]

    for title, body, image in specs:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(6), Inches(2))
        box.text_frame.text = body
        if image:
            slide.shapes.add_picture(
                _swatch(*image), Inches(7.4), Inches(1.6), width=Inches(4.5)
            )

    prs.save(str(target))
    return str(target)


if __name__ == "__main__":
    print(f"Saved {create_sample_presentation()} successfully!")
    print(f"Saved {create_image_presentation()} successfully!")
