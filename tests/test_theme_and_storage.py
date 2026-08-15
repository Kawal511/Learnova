"""
Tests for custom palettes/fonts, Clerk auth guards, and the per-user deck library.
"""

from __future__ import annotations

import pytest

from learnova.auth.clerk import AuthError, user_id_from_header, verify_token
from learnova.rendering.theme_engine import (
    DEFAULT_FONT_ID,
    FONT_CHOICES,
    THEMES,
    apply_font,
    build_custom_theme,
    readable_text_hex,
    resolve_theme,
)


# ── Palette ───────────────────────────────────────────────────────────────────
class TestCustomPalette:
    def test_primary_and_secondary_are_honoured(self):
        theme = build_custom_theme("#1a1a2e", "#e94560")
        assert theme.primary_hex == "#1a1a2e"
        assert theme.accent_hex == "#e94560"
        assert theme.primary_rgb[0] == 0x1A

    def test_hex_without_hash_is_accepted(self):
        assert build_custom_theme("1a1a2e", "e94560").primary_hex == "#1a1a2e"

    def test_three_digit_shorthand_expands(self):
        assert build_custom_theme("#abc", "#def").primary_hex == "#aabbcc"

    def test_garbage_falls_back_instead_of_crashing(self):
        theme = build_custom_theme("not-a-colour", "")
        assert theme.primary_hex.startswith("#")
        assert theme.accent_hex.startswith("#")

    def test_text_is_dark_on_light_background(self):
        assert readable_text_hex("#ffffff") == "#111111"

    def test_text_is_light_on_dark_background(self):
        assert readable_text_hex("#0f172a") == "#ffffff"

    def test_dark_background_produces_readable_body_text(self):
        """A dark canvas must not end up with near-black body text."""
        theme = build_custom_theme("#000000", "#ccff00", background_hex="#111111")
        assert theme.text_hex == "#ffffff"


class TestFonts:
    def test_every_choice_declares_heading_and_body(self):
        for key, value in FONT_CHOICES.items():
            assert value["heading"] and value["body"], key
            assert value["label"]

    def test_custom_theme_uses_selected_pairing(self):
        theme = build_custom_theme("#000", "#fff", font_id="anton_poppins")
        assert theme.heading_font == "Anton"
        assert theme.body_font == "Poppins"

    def test_unknown_font_id_falls_back_to_default(self):
        theme = build_custom_theme("#000", "#fff", font_id="does-not-exist")
        assert theme.heading_font == FONT_CHOICES[DEFAULT_FONT_ID]["heading"]

    def test_apply_font_refonts_a_preset_without_touching_colours(self):
        base = THEMES["swiss_corporate"]
        refonted = apply_font(base, "oswald_lato")
        assert refonted.heading_font == "Oswald"
        assert refonted.primary_hex == base.primary_hex


class TestResolveTheme:
    def test_spec_overrides_named_theme(self):
        theme = resolve_theme("Anything", "midnight_cyber",
                              {"primary": "#123456", "secondary": "#abcdef"})
        assert theme.id == "custom"
        assert theme.primary_hex == "#123456"

    def test_named_theme_used_when_no_spec(self):
        assert resolve_theme("x", "swiss_corporate").id == "swiss_corporate"

    def test_font_only_spec_keeps_preset_colours(self):
        theme = resolve_theme("x", "swiss_corporate", {"font_id": "montserrat"})
        assert theme.id == "swiss_corporate"
        assert theme.heading_font == "Montserrat"

    def test_auto_detection_still_works(self):
        assert resolve_theme("Cyber Security Basics", "auto").id == "midnight_cyber"


class TestThemeReachesArtifacts:
    def test_pptx_embeds_chosen_fonts_and_colours(self):
        import io

        from pptx import Presentation

        from learnova.rendering.ppt_builder import build_pptx

        deck = [{"improved": {"layout_type": "MINIMAL_TEXT", "title": "T",
                              "bullets": ["a", "b"], "takeaway": "t"}}]
        data = build_pptx(deck, "Demo",
                          theme_spec={"primary": "#1a1a2e", "secondary": "#e94560",
                                      "font_id": "anton_poppins"})
        prs = Presentation(io.BytesIO(data))
        names = {
            run.font.name
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
            if run.font.name
        }
        assert "Anton" in names and "Poppins" in names

    def test_web_deck_loads_and_uses_chosen_fonts(self):
        from learnova.rendering.web_deck_builder import build_web_deck

        deck = [{"improved": {"layout_type": "MINIMAL_TEXT", "title": "T",
                              "bullets": ["a"], "takeaway": "t"}}]
        html = build_web_deck(deck, "Demo",
                              theme_spec={"primary": "#111111", "secondary": "#e94560",
                                          "font_id": "oswald_lato"})
        assert "family=Oswald" in html and "family=Lato" in html
        assert "#e94560" in html


# ── Auth guards ───────────────────────────────────────────────────────────────
class TestAuthGuards:
    @pytest.mark.parametrize(
        "header",
        [None, "", "Bearer", "Basic abc", "Bearer ", "token-without-scheme"],
    )
    def test_malformed_authorization_headers_are_rejected(self, header):
        with pytest.raises(AuthError):
            user_id_from_header(header)

    @pytest.mark.parametrize("token", ["", "not-a-jwt", "aaa.bbb.ccc"])
    def test_bad_tokens_are_rejected(self, token):
        with pytest.raises(AuthError):
            verify_token(token)

    def test_unsigned_token_is_not_trusted(self):
        """An alg=none token must never be accepted."""
        import base64
        import json

        def b64(obj):
            return base64.urlsafe_b64encode(json.dumps(obj).encode()).rstrip(b"=").decode()

        forged = f"{b64({'alg': 'none', 'typ': 'JWT'})}.{b64({'sub': 'user_evil'})}."
        with pytest.raises(AuthError):
            verify_token(forged)


# ── Deck library ──────────────────────────────────────────────────────────────
class _FakeStage:
    def __init__(self):
        self.name, self.status, self.seconds = "convert", "ok", 0.1


class _FakeResult:
    def __init__(self):
        self.source_name = "Demo Deck"
        self.markdown = "## Demo\n- a\n"
        self.final_deck = [{"improved": {}}, {"improved": {}}]
        self.quizzes = [{"question": "q"}]
        self.scores = {"overall_score": 77}
        self.converter = "anydoc"
        self.pptx_bytes = b"PK\x03\x04fake"
        self.html_bytes = b"<html></html>"
        self.stages = [_FakeStage()]


@pytest.fixture
def library(tmp_path, monkeypatch):
    """Point the library at a temp dir so tests never touch real user data."""
    from learnova.storage import deck_library as lib

    monkeypatch.setattr(lib, "DATA_DIR", tmp_path)
    return lib


class TestDeckLibrary:
    def test_save_then_list_round_trips(self, library):
        record = library.save_deck("user_alpha", _FakeResult(), theme_id="auto")
        decks = library.list_decks("user_alpha")
        assert len(decks) == 1
        assert decks[0]["id"] == record.id
        assert decks[0]["slide_count"] == 2
        assert decks[0]["overall_score"] == 77

    def test_users_cannot_see_each_others_decks(self, library):
        library.save_deck("user_alpha", _FakeResult())
        assert library.list_decks("user_beta") == []

    def test_artifacts_and_markdown_are_retrievable(self, library):
        record = library.save_deck("user_alpha", _FakeResult())
        assert library.read_artifact("user_alpha", record.id, "pptx").startswith(b"PK")
        assert library.read_artifact("user_alpha", record.id, "html") == b"<html></html>"
        assert "## Demo" in library.read_markdown("user_alpha", record.id)

    def test_another_user_cannot_read_artifacts(self, library):
        record = library.save_deck("user_alpha", _FakeResult())
        assert library.read_artifact("user_beta", record.id, "pptx") is None
        assert library.read_markdown("user_beta", record.id) is None

    def test_theme_spec_is_persisted_for_display(self, library):
        spec = {"primary": "#1a1a2e", "secondary": "#e94560"}
        record = library.save_deck("user_alpha", _FakeResult(), theme_spec=spec)
        assert library.get_deck("user_alpha", record.id)["theme_spec"] == spec

    def test_delete_removes_the_deck(self, library):
        record = library.save_deck("user_alpha", _FakeResult())
        assert library.delete_deck("user_alpha", record.id) is True
        assert library.list_decks("user_alpha") == []
        assert library.delete_deck("user_alpha", record.id) is False

    def test_newest_deck_is_listed_first(self, library):
        first = library.save_deck("user_alpha", _FakeResult(), title="Older")
        second = library.save_deck("user_alpha", _FakeResult(), title="Newer")
        titles = [d["title"] for d in library.list_decks("user_alpha")]
        assert titles.index("Newer") < titles.index("Older")
        assert first.id != second.id

    @pytest.mark.parametrize("evil", ["../escape", "a/b", "..", "user id", "x" * 200, ""])
    def test_path_traversal_is_rejected(self, library, evil):
        """A crafted id must never let one user reach outside their directory."""
        with pytest.raises(ValueError):
            library.save_deck(evil, _FakeResult())

    def test_traversal_in_deck_id_is_rejected(self, library):
        library.save_deck("user_alpha", _FakeResult())
        with pytest.raises(ValueError):
            library.read_artifact("user_alpha", "../../etc", "pptx")

    def test_listing_an_unknown_user_is_empty_not_an_error(self, library):
        assert library.list_decks("user_nobody") == []
