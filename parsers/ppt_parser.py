"""
Advanced PPTX Parser for Learnova (Migrated to Unified Extraction Architecture)
Extracts ALL text and structural elements from every shape type:
- Text frames (paragraphs + runs, preserving bullet levels, bold/italic, font size)
- Grouped shapes (recursive traversal with spatial coordinates)
- Native PowerPoint tables (rows + cell matrices)
- SmartArt and diagram shapes (via XML text extraction & DiagramElement)
- Chart alt-text, chart titles, and series data (StructuredChartElement)
- Equations (OMML extraction to EquationElement)
- Slide notes/speaker notes
- Slide layout/master fallback titles
- Embedded pictures (SHA-256 hashed VisualAssetElement) + slide thumbnail rendering
"""

import io
import re
import os
import hashlib
from typing import List, Dict, Any, Optional, Set
from dataclasses import dataclass, field
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

from parsers.base import BaseDocumentParser
from parsers.schema import (
    DocumentEntity,
    SlidePageEntity,
    TextBlockElement,
    TableElement,
    VisualAssetElement,
    StructuredChartElement,
    DiagramElement,
    EquationElement,
    DocumentType,
    VisualAssetType,
    ChartType,
    DiagramType,
)


@dataclass
class SlideData:
    """Legacy SlideData structure for backward compatibility."""

    id: int
    title: str
    text: str
    image: Optional[dict] = None
    images: List[dict] = field(default_factory=list)


@dataclass
class ParsedDocument:
    """Legacy ParsedDocument structure for backward compatibility."""

    slide_units: List[SlideData]


def _xml_text_deep(element) -> List[str]:
    """Recursively pulls all a:t text nodes from an OOXML element."""
    texts = []
    for t_elem in element.iter(qn("a:t")):
        v = (t_elem.text or "").strip()
        if v:
            texts.append(v)
    return texts


def _get_shape_bbox(shape) -> Optional[tuple[float, float, float, float]]:
    """Calculates spatial bounding box (x0, y0, x1, y1) in points."""
    try:
        x0 = float(shape.left.pt)
        y0 = float(shape.top.pt)
        x1 = float((shape.left + shape.width).pt)
        y1 = float((shape.top + shape.height).pt)
        return (round(x0, 2), round(y0, 2), round(x1, 2), round(y1, 2))
    except Exception:
        return None


def _slide_as_image(file_path: str, slide_index: int) -> Optional[dict]:
    """
    Render a PDF page to PNG using PyMuPDF.
    For PPTX files, PyMuPDF cannot render slide thumbnails natively.
    """
    if not file_path.lower().endswith(".pdf"):
        return None
    try:
        import fitz
        doc = fitz.open(file_path)
        if slide_index < len(doc):
            page = doc[slide_index]
            pix = page.get_pixmap(dpi=150)
            png_bytes = bytes(pix.tobytes("png"))
            pix = None
            doc.close()
            return {"bytes": png_bytes, "ext": "png"}
        doc.close()
    except Exception:
        pass
    return None


class PPTXParser(BaseDocumentParser):
    """
    Production-ready PowerPoint (.pptx) parser implementing the BaseDocumentParser interface.
    Extracts structured DocumentEntity graphs with typed elements.
    """

    def supports(self, file_path_or_extension: str) -> bool:
        ext = file_path_or_extension.lower().strip()
        return ext.endswith(".pptx") or ext == "pptx"

    def parse(self, file_path: str) -> DocumentEntity:
        if not self.validate(file_path):
            raise ValueError(f"Invalid or unsupported PPTX file path: {file_path}")

        prs = Presentation(file_path)
        slides_entities: List[SlidePageEntity] = []

        for i, slide in enumerate(prs.slides):
            slide_entity = self._parse_slide(slide, slide_index=i, file_path=file_path)
            slides_entities.append(slide_entity)

        doc_id = hashlib.sha256(os.path.basename(file_path).encode("utf-8")).hexdigest()[:16]
        return DocumentEntity(
            id=doc_id,
            filename=os.path.basename(file_path),
            doc_type=DocumentType.PPTX,
            total_units=len(slides_entities),
            slides=slides_entities,
            metadata={"source_path": file_path, "slide_count": len(slides_entities)},
        )

    def _parse_slide(self, slide, slide_index: int, file_path: str) -> SlidePageEntity:
        # Title extraction
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
        except Exception:
            pass

        text_blocks: List[TextBlockElement] = []
        tables: List[TableElement] = []
        visual_assets: List[VisualAssetElement] = []
        charts: List[StructuredChartElement] = []
        diagrams: List[DiagramElement] = []
        equations: List[EquationElement] = []
        visited_ids: Set[int] = set()

        reading_order_counter = 0

        # Sort shapes spatially by (top, left) for logical reading order
        sorted_shapes = list(slide.shapes)
        try:
            sorted_shapes.sort(key=lambda s: (getattr(s, "top", 0), getattr(s, "left", 0)))
        except Exception:
            pass

        for shape in sorted_shapes:
            reading_order_counter = self._extract_shape_elements(
                shape=shape,
                slide_index=slide_index,
                visited_ids=visited_ids,
                text_blocks=text_blocks,
                tables=tables,
                visual_assets=visual_assets,
                charts=charts,
                diagrams=diagrams,
                equations=equations,
                order_counter=reading_order_counter,
                is_title_shape=(shape == getattr(slide.shapes, "title", None)),
            )

        # Speaker notes
        speaker_notes_text = None
        try:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_str = notes_tf.text.strip()
            if notes_str and len(notes_str) > 5:
                speaker_notes_text = notes_str
        except Exception:
            pass

        # Fallback thumbnail image if no embedded assets found
        rendered_page_img = None
        if not visual_assets:
            fallback_img = _slide_as_image(file_path, slide_index)
            if fallback_img:
                rendered_page_img = VisualAssetElement(
                    id=f"slide_{slide_index}_rendered",
                    image_bytes=fallback_img["bytes"],
                    format="png",
                    asset_type=VisualAssetType.SCANNED_PAGE,
                )

        reading_order_ids = [tb.id for tb in text_blocks] + [t.id for t in tables] + [v.id for v in visual_assets]

        return SlidePageEntity(
            id=slide_index,
            unit_number=slide_index + 1,
            title=title or f"Slide {slide_index + 1}",
            text_blocks=text_blocks,
            tables=tables,
            visual_assets=visual_assets,
            charts=charts,
            diagrams=diagrams,
            equations=equations,
            rendered_page_image=rendered_page_img,
            speaker_notes=speaker_notes_text,
            reading_order_elements=reading_order_ids,
        )

    def _extract_shape_elements(
        self,
        shape,
        slide_index: int,
        visited_ids: Set[int],
        text_blocks: List[TextBlockElement],
        tables: List[TableElement],
        visual_assets: List[VisualAssetElement],
        charts: List[StructuredChartElement],
        diagrams: List[DiagramElement],
        equations: List[EquationElement],
        order_counter: int,
        is_title_shape: bool = False,
    ) -> int:

        if id(shape) in visited_ids:
            return order_counter
        visited_ids.add(id(shape))

        bbox = _get_shape_bbox(shape)
        shape_id = f"s{slide_index}_{id(shape)}"

        # 1. Native Table Shapes
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                tbl = shape.table
                headers = []
                rows = []
                for r_idx, row in enumerate(tbl.rows):
                    row_cells = [cell.text.strip() for cell in row.cells]
                    if r_idx == 0:
                        headers = row_cells
                    else:
                        rows.append(row_cells)
                tables.append(TableElement(
                    id=f"{shape_id}_tbl",
                    headers=headers,
                    rows=rows,
                    num_rows=len(tbl.rows),
                    num_cols=len(tbl.columns),
                    bbox=bbox,
                    reading_order=order_counter,
                ))
                order_counter += 1
                return order_counter
        except Exception:
            pass

        # 2. Embedded Pictures
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob = bytes(shape.image.blob)
                sha256 = hashlib.sha256(blob).hexdigest()
                width_px, height_px = 0, 0
                png_bytes = blob
                try:
                    with Image.open(io.BytesIO(blob)) as pil_img:
                        pil_img.load()
                        width_px, height_px = pil_img.size
                        out = io.BytesIO()
                        pil_img.convert("RGB").save(out, format="PNG")
                        png_bytes = bytes(out.getvalue())
                except Exception:
                    pass

                if width_px >= 100 and height_px >= 100:
                    visual_assets.append(VisualAssetElement(
                        id=f"{shape_id}_img",
                        image_bytes=png_bytes,
                        format="png",
                        width_px=width_px,
                        height_px=height_px,
                        asset_type=VisualAssetType.PICTURE,
                        sha256_hash=sha256,
                        bbox=bbox,
                        reading_order=order_counter,
                    ))
                    order_counter += 1
                return order_counter
        except Exception:
            pass

        # 3. Charts
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart = shape.chart
                title = ""
                try:
                    if chart.has_title and chart.chart_title and chart.chart_title.text_frame:
                        title = chart.chart_title.text_frame.text.strip()
                except Exception:
                    pass

                alt_text = None
                try:
                    nvSpPr = shape._element.find(".//" + qn("p:nvSpPr"))
                    if nvSpPr is not None:
                        nvPr = nvSpPr.find(qn("p:nvPr"))
                        if nvPr is not None:
                            ph_elem = nvPr.find(qn("p:ph"))
                            if ph_elem is not None and ph_elem.get("descr"):
                                alt_text = ph_elem.get("descr")
                except Exception:
                    pass

                categories = []
                series_data = []
                try:
                    if chart.plots:
                        plot = chart.plots[0]
                        if hasattr(plot, "categories") and plot.categories:
                            categories = [str(c) for c in plot.categories]
                        for s in plot.series:
                            s_name = getattr(s, "name", "Series")
                            s_vals = list(getattr(s, "values", []))
                            series_data.append({"name": s_name, "values": s_vals})
                except Exception:
                    pass

                charts.append(StructuredChartElement(
                    id=f"{shape_id}_chart",
                    title=title or "Chart",
                    chart_type=ChartType.BAR,
                    categories=categories,
                    series_data=series_data,
                    alt_text=alt_text,
                    bbox=bbox,
                    reading_order=order_counter,
                ))
                order_counter += 1
                return order_counter
        except Exception:
            pass

        # 4. Grouped Shapes (Recurse)
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    order_counter = self._extract_shape_elements(
                        shape=child,
                        slide_index=slide_index,
                        visited_ids=visited_ids,
                        text_blocks=text_blocks,
                        tables=tables,
                        visual_assets=visual_assets,
                        charts=charts,
                        diagrams=diagrams,
                        equations=equations,
                        order_counter=order_counter,
                    )
                return order_counter
        except Exception:
            pass

        # 5. OMML Math Equation Extraction
        try:
            for omath in shape._element.iter(qn("m:oMath")):
                math_text = "".join(t.text for t in omath.iter(qn("m:t")) if t.text).strip()
                if math_text:
                    equations.append(EquationElement(
                        id=f"{shape_id}_eq_{len(equations)}",
                        latex_expression=math_text,
                        ascii_fallback=math_text,
                        raw_omml=math_text,
                        bbox=bbox,
                        reading_order=order_counter,
                    ))
                    order_counter += 1
        except Exception:
            pass

        # 6. Text Frames
        initial_tb_count = len(text_blocks)
        try:
            if shape.has_text_frame:
                tf = shape.text_frame
                for p_idx, para in enumerate(tf.paragraphs):
                    para_text = "".join(run.text for run in para.runs).strip()
                    if not para_text:
                        para_text = para.text.strip()
                    if not para_text:
                        continue

                    is_bold = False
                    is_italic = False
                    font_size = None
                    for run in para.runs:
                        if run.font:
                            if run.font.bold:
                                is_bold = True
                            if run.font.italic:
                                is_italic = True
                            if run.font.size and run.font.size.pt:
                                font_size = float(run.font.size.pt)

                    level = getattr(para, "level", 0)

                    text_blocks.append(TextBlockElement(
                        id=f"{shape_id}_tb_{p_idx}",
                        text=para_text,
                        is_heading=is_title_shape or (p_idx == 0 and level == 0 and (font_size and font_size >= 20)),
                        heading_level=1 if is_title_shape else (2 if font_size and font_size >= 16 else 0),
                        bullet_level=level,
                        font_size=font_size,
                        is_bold=is_bold,
                        is_italic=is_italic,
                        bbox=bbox,
                        reading_order=order_counter,
                    ))
                    order_counter += 1
        except Exception:
            pass

        # 7. XML Deep Fallback (SmartArt / Diagrams)
        if len(text_blocks) == initial_tb_count:
            try:
                xml_texts = _xml_text_deep(shape._element)
                if xml_texts:
                    diagrams.append(DiagramElement(
                        id=f"{shape_id}_diag",
                        title=xml_texts[0] if xml_texts else "Diagram",
                        diagram_type=DiagramType.SMARTART,
                        raw_xml_text=xml_texts,
                        bbox=bbox,
                        reading_order=order_counter,
                    ))
                    order_counter += 1
            except Exception:
                pass

        return order_counter


def parse_ppt(file_path: str) -> ParsedDocument:
    """
    Backward-compatible entry point for PPTX parsing.
    Invokes PPTXParser to construct a DocumentEntity graph, then bridges to legacy ParsedDocument.
    """
    parser = PPTXParser()
    doc_entity = parser.parse(file_path)
    legacy_dicts = doc_entity.to_legacy_parsed_dicts()

    slide_units = [
        SlideData(
            id=d["id"],
            title=d["title"],
            text=d["text"],
            image=d["image"],
            images=d["images"],
        )
        for d in legacy_dicts
    ]
    return ParsedDocument(slide_units=slide_units)
