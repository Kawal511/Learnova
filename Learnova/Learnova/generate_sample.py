from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)

# Slide 1: Title
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]
title1.text = "Artificial Intelligence in Education"
subtitle1.text = "Transforming Outdated Courseware into Engaging Learning Modules"

# Slide 2: Minimal text list
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title2 = slide2.shapes.title
title2.text = "Key Challenges in Higher Education"
tf2 = slide2.placeholders[1].text_frame
tf2.text = "Traditional presentations suffer from heavy text density and passive student engagement."
p1 = tf2.add_paragraph()
p1.text = "Lack of interactive assessment checkpoints during lectures reduces retention rates."
p2 = tf2.add_paragraph()
p2.text = "Manual redesign of slides requires hours of design expertise per lecture."

# Slide 3: Process workflow (flowchart candidate)
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
title3.text = "Learnova Automated Processing Pipeline Workflow"
tf3 = slide3.placeholders[1].text_frame
tf3.text = "Step 1: Upload raw PPTX or PDF textbook."
p31 = tf3.add_paragraph()
p31.text = "Step 2: Parse document layout, tables, SmartArt, and embedded diagram images."
p32 = tf3.add_paragraph()
p32.text = "Step 3: Route content through AI Layout Router to classify dynamic visual structures."
p33 = tf3.add_paragraph()
p33.text = "Step 4: Generate interactive quizzes and export animated PPTX and HTML web decks."

# Slide 4: Metric callout candidate
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title4 = slide4.shapes.title
title4.text = "Student Engagement Benchmark Metrics"
tf4 = slide4.placeholders[1].text_frame
tf4.text = "Studies show 84% improvement in concept retention when interactive checkpoints and dynamic visual layouts are integrated into lecture slides."

# Slide 5: Card Grid candidate
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title5 = slide5.shapes.title
title5.text = "Four Pillars of Learnova Architecture"
tf5 = slide5.placeholders[1].text_frame
tf5.text = "1. AI Visual Engine: Converts plain text to flowcharts and structured cards."
p51 = tf5.add_paragraph()
p51.text = "2. Vision OCR Integration: Reads embedded diagrams and infographic text."
p52 = tf5.add_paragraph()
p52.text = "3. Interleaved Quizzes: Promotes active recall with checkpoint questions."
p53 = tf5.add_paragraph()
p53.text = "4. Multi-Format Export: Native PPTX presentation + 60fps web deck."

prs.save("/Users/swayampanchal/Desktop/New Learnova/Learnova/sample_test_presentation.pptx")
print("Saved sample_test_presentation.pptx successfully!")
