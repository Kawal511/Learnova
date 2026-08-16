"""
Tests for dynamic slide geometry, the inline quiz band, and the layout
classification tightening.
"""

from __future__ import annotations

import io

import pytest
from pptx import Presentation
from pptx.util import Emu

from learnova.ai.layout_router import (
    _build_fallback,
    _heuristic_layout_type,
    parse_pipe_table,
    split_sentences,
)
from learnova.ai.quiz_gen import interleave_quizzes_into_slides
from learnova.rendering import layout as L
from learnova.rendering.ppt_builder import build_pptx

SLIDE_BOTTOM = 7.5
SLIDE_RIGHT = 13.33


# ── Geometry ──────────────────────────────────────────────────────────────────
class TestContentBand:
    def test_band_shrinks_for_takeaway(self):
        plain = L.content_band(has_takeaway=False, has_quiz=False)
        with_tk = L.content_band(has_takeaway=True, has_quiz=False)
        assert with_tk.height < plain.height

    def test_band_shrinks_further_for_quiz(self):
        with_tk = L.content_band(has_takeaway=True, has_quiz=False)
        with_both = L.content_band(has_takeaway=True, has_quiz=True)
        assert with_both.height < with_tk.height

    def test_bands_never_leave_the_canvas(self):
        for tk in (True, False):
            for qz in (True, False):
                band = L.content_band(tk, qz)
                assert band.top >= L.HEADER_H
                assert band.top + band.height <= SLIDE_BOTTOM
                assert band.height > 0

    def test_quiz_band_sits_above_the_takeaway(self):
        quiz = L.quiz_band(has_takeaway=True)
        takeaway = L.takeaway_band()
        assert quiz.top + quiz.height <= takeaway.top + 0.01

    def test_body_and_quiz_do_not_overlap(self):
        band = L.content_band(has_takeaway=True, has_quiz=True)
        quiz = L.quiz_band(has_takeaway=True)
        assert band.top + band.height <= quiz.top + 0.01


class TestFontFitting:
    def test_short_text_gets_the_maximum_size(self):
        assert L.fit_font_size(["Hi"], 10.0, 4.0) == L.MAX_BODY_PT

    def test_long_text_shrinks(self):
        long_text = ["word " * 400]
        assert L.fit_font_size(long_text, 6.0, 2.0) < L.MAX_BODY_PT

    def test_never_below_the_legibility_floor(self):
        assert L.fit_font_size(["x " * 5000], 3.0, 0.5) == L.MIN_BODY_PT

    def test_more_text_never_gets_a_bigger_size(self):
        small = L.fit_font_size(["one line"], 8.0, 3.0)
        big = L.fit_font_size(["one line " * 60], 8.0, 3.0)
        assert big <= small

    def test_empty_input_is_safe(self):
        assert L.fit_font_size([], 8.0, 3.0) == L.MAX_BODY_PT
        assert L.fit_font_size(["", "  "], 8.0, 3.0) == L.MAX_BODY_PT

    def test_estimate_lines_grows_with_length(self):
        assert L.estimate_lines("x" * 400, 5.0, 14) > L.estimate_lines("x" * 40, 5.0, 14)


class TestGridCells:
    @pytest.mark.parametrize("count", [1, 2, 3, 4, 5, 6, 8])
    def test_cells_stay_inside_the_area(self, count):
        area = L.Box(0.5, 1.35, 12.33, 4.0)
        for cell in L.grid_cells(count, area):
            assert cell.left >= area.left - 0.001
            assert cell.top >= area.top - 0.001
            assert cell.left + cell.width <= area.left + area.width + 0.001
            assert cell.top + cell.height <= area.top + area.height + 0.001

    def test_returns_one_cell_per_item(self):
        assert len(L.grid_cells(5, L.Box(0, 0, 12, 4))) == 5

    def test_wraps_past_four_per_row(self):
        cells = L.grid_cells(6, L.Box(0, 0, 12, 4))
        assert len({round(c.top, 3) for c in cells}) == 2

    def test_zero_items_is_empty(self):
        assert L.grid_cells(0, L.Box(0, 0, 12, 4)) == []

    def test_split_text_image_partitions_without_overlap(self):
        area = L.Box(0.5, 1.35, 12.33, 4.0)
        text, image = L.split_text_image(area, has_image=True)
        assert text.left + text.width <= image.left + 0.001
        assert image.left + image.width <= area.left + area.width + 0.001

    def test_no_image_gives_full_width(self):
        area = L.Box(0.5, 1.35, 12.33, 4.0)
        text, image = L.split_text_image(area, has_image=False)
        assert text.width == area.width
        assert image.width == 0


# ── Inline quiz ───────────────────────────────────────────────────────────────
def _deck(n=8):
    return [{"original": {"text": "x"}, "improved": {
        "layout_type": "MINIMAL_TEXT", "title": f"Topic {i}",
        "bullets": ["A point about the topic", "Another point"],
        "takeaway": "Remember this."}} for i in range(1, n + 1)]


QUIZ = {
    "question": "Which statistic does ANOVA compare?",
    "options": ["A) Variance", "B) Median", "C) Mode", "D) Range"],
    "correct": "A",
    "explanation": "It compares variances.",
}


class TestInlineQuiz:
    def test_inline_adds_no_slides(self):
        deck = _deck(8)
        out = interleave_quizzes_into_slides(deck, [QUIZ, QUIZ], frequency=4)
        assert len(out) == len(deck)

    def test_attached_to_every_nth_slide(self):
        out = interleave_quizzes_into_slides(_deck(8), [QUIZ, QUIZ], frequency=4)
        carriers = [i for i, e in enumerate(out) if e["improved"].get("inline_quiz")]
        assert carriers == [3, 7]

    def test_original_slide_content_is_preserved(self):
        out = interleave_quizzes_into_slides(_deck(4), [QUIZ], frequency=4)
        assert out[3]["improved"]["bullets"]
        assert out[3]["improved"]["layout_type"] == "MINIMAL_TEXT"

    def test_options_are_capped_at_four(self):
        many = {**QUIZ, "options": [f"Opt {i}" for i in range(9)]}
        out = interleave_quizzes_into_slides(_deck(4), [many], frequency=4)
        assert len(out[3]["improved"]["inline_quiz"]["options"]) == 4

    def test_separate_slide_mode_still_available(self):
        deck = _deck(4)
        out = interleave_quizzes_into_slides(deck, [QUIZ], frequency=4, inline=False)
        assert len(out) == len(deck) + 1
        assert out[-1]["improved"]["layout_type"] == "QUIZ"

    def test_no_quizzes_returns_the_deck_untouched(self):
        deck = _deck(4)
        assert interleave_quizzes_into_slides(deck, [], frequency=4) is deck

    def test_rendered_band_fits_on_the_slide(self):
        out = interleave_quizzes_into_slides(_deck(4), [QUIZ], frequency=4)
        prs = Presentation(io.BytesIO(build_pptx(out, "T", theme_id="brutalist_neon")))
        slide = prs.slides[4]                       # +1 for the title slide
        for shape in slide.shapes:
            assert Emu(shape.top).inches + Emu(shape.height).inches <= SLIDE_BOTTOM + 0.05
            assert Emu(shape.left).inches + Emu(shape.width).inches <= SLIDE_RIGHT + 0.05

    def test_option_letters_are_not_doubled(self):
        out = interleave_quizzes_into_slides(_deck(4), [QUIZ], frequency=4)
        prs = Presentation(io.BytesIO(build_pptx(out, "T", theme_id="brutalist_neon")))
        texts = [sh.text_frame.text for sh in prs.slides[4].shapes if sh.has_text_frame]
        assert any(t.strip() == "A. Variance" for t in texts), texts


# ── Layout classification ─────────────────────────────────────────────────────
class TestTableClassification:
    def test_real_pipe_table_is_detected(self):
        headers, rows = parse_pipe_table("|W|M|\n|---|---|\n|1|18|\n|2|19|")
        assert headers == ["W", "M"]
        assert rows == [["1", "18"], ["2", "19"]]

    def test_prose_has_no_table(self):
        assert parse_pipe_table("Just a sentence about tables.") == ([], [])

    def test_single_row_is_not_a_table(self):
        assert parse_pipe_table("|only|one|")[1] == []

    @pytest.mark.parametrize("text", [
        "See the table below for details.",
        "Advantages of this approach are many.",
        "The value of vsomething is unclear.",
    ])
    def test_stray_keywords_no_longer_force_a_table(self, text):
        assert _heuristic_layout_type(text) != "TABLE"

    @pytest.mark.parametrize("text", [
        "Method A vs Method B differ.",
        "A comparison of the two designs.",
        "|a|b|\n|---|---|\n|1|2|\n|3|4|",
    ])
    def test_genuine_signals_still_give_a_table(self, text):
        assert _heuristic_layout_type(text) == "TABLE"

    def test_table_without_rows_downgrades_instead_of_faking_one(self):
        result = _build_fallback("Method A vs Method B differ.", "T", "TABLE")
        assert result["layout_type"] == "MINIMAL_TEXT"
        assert "table_rows" not in result

    def test_table_with_rows_uses_the_real_data(self):
        result = _build_fallback("|W|M|\n|---|---|\n|1|18|\n|2|19|", "T", "TABLE")
        assert result["layout_type"] == "TABLE"
        assert result["table_rows"] == [["1", "18"], ["2", "19"]]


class TestSentenceSplitting:
    def test_abbreviation_is_not_a_break(self):
        assert len(split_sentences("Total no. of observations is twelve.")) == 1

    def test_decimal_is_not_a_break(self):
        assert len(split_sentences("The mean is 3.14 exactly.")) == 1

    def test_real_boundaries_still_split(self):
        assert len(split_sentences("First one. Second one. Third one.")) == 3

    def test_empty_is_safe(self):
        assert split_sentences("") == []
